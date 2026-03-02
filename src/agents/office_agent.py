"""Office format agents: DocxAgent, PptxAgent, XlsxAgent."""
import logging
from pathlib import Path

logger = logging.getLogger(__name__)

# Cap xlsx output so we don't index enormous spreadsheets verbatim
_XLSX_MAX_ROWS = 200
_XLSX_MAX_COLS = 20


class DocxAgent:
    """Convert .docx files to Markdown, including tables and embedded images."""

    def __init__(self, image_agent, images_dir: Path):
        self.image_agent = image_agent
        self.images_dir = images_dir
        self.images_dir.mkdir(parents=True, exist_ok=True)

    def process(self, file_path: Path) -> dict:
        from docx import Document
        from docx.table import Table
        from docx.text.paragraph import Paragraph as DocxParagraph

        doc = Document(str(file_path))
        metadata = self._get_metadata(file_path, doc)

        header = f"# {metadata['title']}\n"
        if metadata.get("author"):
            header += f"\n**Author:** {metadata['author']}\n"
        header += "\n---"

        parts = [header]

        # Iterate paragraphs and tables in document order via the XML body
        for child in doc.element.body:
            tag = child.tag.split("}")[-1] if "}" in child.tag else child.tag
            if tag == "p":
                para = DocxParagraph(child, doc)
                md = self._para_to_md(para)
                if md:
                    parts.append(md)
            elif tag == "tbl":
                table = Table(child, doc)
                md = self._table_to_markdown(table)
                if md:
                    parts.append(md)

        # Embedded images (position within doc is not easily recoverable, append at end)
        img_count = 0
        for rel in doc.part.rels.values():
            if "image" not in rel.reltype:
                continue
            try:
                img_bytes = rel.target_part.blob
                content_type = rel.target_part.content_type  # e.g. "image/png"
                ext = content_type.split("/")[-1].replace("jpeg", "jpg")
                img_path = self.images_dir / f"{file_path.stem}_img{img_count}.{ext}"
                img_path.write_bytes(img_bytes)
                desc = self.image_agent.describe(
                    img_path, f"a Word document titled '{metadata['title']}'"
                )
                parts.append(desc)
                img_count += 1
            except Exception as e:
                logger.warning("docx image %d skipped: %s", img_count, e)

        return {"markdown": "\n\n".join(parts), "metadata": metadata}

    def _para_to_md(self, para) -> str:
        text = para.text.strip()
        if not text:
            return ""
        style = para.style.name if para.style else ""
        if style.startswith("Heading 1"):
            return f"# {text}"
        if style.startswith("Heading 2"):
            return f"## {text}"
        if style.startswith("Heading 3"):
            return f"### {text}"
        if style.startswith("Heading 4"):
            return f"#### {text}"
        if "List" in style:
            level = getattr(para, "level", 0) or 0
            indent = "  " * level
            return f"{indent}- {text}"
        return text

    def _table_to_markdown(self, table) -> str:
        rows = table.rows
        if not rows:
            return ""

        def cell_text(cell):
            return cell.text.replace("\n", " ").strip()

        header = [cell_text(c) for c in rows[0].cells]
        # Deduplicate merged cells (python-docx repeats the first cell value)
        seen_ids = set()
        unique_header = []
        for i, c in enumerate(rows[0].cells):
            cid = id(c._tc)
            unique_header.append(cell_text(c) if cid not in seen_ids else "")
            seen_ids.add(cid)
        header = unique_header

        md = "| " + " | ".join(header) + " |\n"
        md += "| " + " | ".join("---" for _ in header) + " |\n"
        for row in rows[1:]:
            seen_ids = set()
            cells = []
            for c in row.cells:
                cid = id(c._tc)
                cells.append(cell_text(c) if cid not in seen_ids else "")
                seen_ids.add(cid)
            md += "| " + " | ".join(cells) + " |\n"
        return md

    def _get_metadata(self, file_path: Path, doc) -> dict:
        cp = doc.core_properties
        return {
            "source_file": str(file_path),
            "file_type": "docx",
            "title": cp.title or file_path.stem,
            "author": cp.author or "",
            "creation_date": str(cp.created or ""),
        }


class PptxAgent:
    """Convert .pptx files to Markdown, one section per slide."""

    def __init__(self, image_agent, images_dir: Path):
        self.image_agent = image_agent
        self.images_dir = images_dir
        self.images_dir.mkdir(parents=True, exist_ok=True)

    def process(self, file_path: Path) -> dict:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        metadata = self._get_metadata(file_path, prs)

        slide_sections = [f"# {metadata['title']}\n\n---"]
        img_count = 0

        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_parts = [f"## Slide {slide_num}"]

            # Title shape first
            title_text = ""
            if slide.shapes.title and slide.shapes.title.has_text_frame:
                title_text = slide.shapes.title.text.strip()
                if title_text:
                    slide_parts.append(f"### {title_text}")

            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue  # already handled

                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if not text:
                            continue
                        indent = "  " * (para.level or 0)
                        slide_parts.append(f"{indent}- {text}" if para.level else text)

                elif hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        ext = shape.image.ext
                        img_path = (
                            self.images_dir
                            / f"{file_path.stem}_slide{slide_num}_img{img_count}.{ext}"
                        )
                        img_path.write_bytes(img_bytes)
                        context = f"slide {slide_num}"
                        if title_text:
                            context += f" titled '{title_text}'"
                        desc = self.image_agent.describe(
                            img_path,
                            f"a PowerPoint presentation titled '{metadata['title']}' ({context})",
                        )
                        slide_parts.append(desc)
                        img_count += 1
                    except Exception as e:
                        logger.warning("slide %d image %d skipped: %s", slide_num, img_count, e)

            # Speaker notes
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_parts.append(f"*Notes: {notes}*")

            slide_sections.append("\n\n".join(slide_parts))

        return {
            "markdown": "\n\n---\n\n".join(slide_sections),
            "metadata": metadata,
        }

    def _get_metadata(self, file_path: Path, prs) -> dict:
        cp = prs.core_properties
        return {
            "source_file": str(file_path),
            "file_type": "pptx",
            "title": cp.title or file_path.stem,
            "author": cp.author or "",
            "creation_date": str(cp.created or ""),
            "slide_count": len(prs.slides),
        }


class XlsxAgent:
    """Convert .xlsx files to Markdown tables, one section per sheet."""

    def process(self, file_path: Path) -> dict:
        import openpyxl

        wb = openpyxl.load_workbook(str(file_path), data_only=True)
        metadata = self._get_metadata(file_path, wb)

        parts = [f"# {metadata['title']}\n\n---"]
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"## Sheet: {sheet_name}")
            parts.append(self._sheet_to_markdown(ws))

        return {"markdown": "\n\n".join(parts), "metadata": metadata}

    def _sheet_to_markdown(self, ws) -> str:
        rows = list(
            ws.iter_rows(
                max_row=_XLSX_MAX_ROWS, max_col=_XLSX_MAX_COLS, values_only=True
            )
        )
        # Drop trailing all-None rows
        while rows and all(v is None for v in rows[-1]):
            rows.pop()
        if not rows:
            return "*Empty sheet*"

        def fmt(v) -> str:
            return "" if v is None else str(v).replace("\n", " ").strip()

        header = [fmt(c) for c in rows[0]]
        md = "| " + " | ".join(header) + " |\n"
        md += "| " + " | ".join("---" for _ in header) + " |\n"
        for row in rows[1:]:
            padded = list(row) + [None] * max(0, len(header) - len(row))
            md += "| " + " | ".join(fmt(c) for c in padded[: len(header)]) + " |\n"
        return md

    def _get_metadata(self, file_path: Path, wb) -> dict:
        cp = wb.properties
        return {
            "source_file": str(file_path),
            "file_type": "xlsx",
            "title": cp.title or file_path.stem,
            "author": cp.creator or "",
            "creation_date": str(cp.created or ""),
        }
