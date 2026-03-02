"""Tests for DocxAgent, PptxAgent, and XlsxAgent.

Images are skipped via a lightweight mock ImageProcessingAgent so these tests
run without hitting the Anthropic API.
"""
from pathlib import Path
from unittest.mock import MagicMock

import pytest


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _mock_image_agent():
    """Returns a mock that records calls and returns a fixed description."""
    mock = MagicMock()
    mock.describe.return_value = "*[image described by Claude Vision]*"
    return mock


def _make_docx(path: Path, *, headings=True, tables=True) -> Path:
    """Write a minimal .docx file using python-docx."""
    from docx import Document
    doc = Document()
    doc.core_properties.title = "Test Document"
    doc.core_properties.author = "Test Author"
    if headings:
        doc.add_heading("Section One", level=1)
        doc.add_heading("Subsection", level=2)
    doc.add_paragraph("This is a body paragraph.")
    if tables:
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Name"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "alpha"
        table.cell(1, 1).text = "1"
    doc.save(str(path))
    return path


def _make_pptx(path: Path) -> Path:
    """Write a minimal .pptx file using python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.core_properties.title = "Test Presentation"
    prs.core_properties.author = "Presenter"

    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Slide Title"
    slide.placeholders[1].text = "Bullet point content"

    prs.save(str(path))
    return path


def _make_xlsx(path: Path) -> Path:
    """Write a minimal .xlsx file using openpyxl."""
    import openpyxl

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Data"
    ws.append(["City", "Population", "Country"])
    ws.append(["Tokyo", 13960000, "Japan"])
    ws.append(["Berlin", 3645000, "Germany"])
    wb.save(str(path))
    return path


# ---------------------------------------------------------------------------
# DocxAgent tests
# ---------------------------------------------------------------------------

class TestDocxAgent:
    @pytest.fixture
    def agent(self, tmp_path):
        from src.agents.office_agent import DocxAgent
        return DocxAgent(image_agent=_mock_image_agent(), images_dir=tmp_path / "images")

    @pytest.fixture
    def docx_file(self, tmp_path):
        return _make_docx(tmp_path / "test.docx")

    def test_returns_markdown_and_metadata(self, agent, docx_file):
        result = agent.process(docx_file)
        assert "markdown" in result
        assert "metadata" in result

    def test_title_in_markdown(self, agent, docx_file):
        result = agent.process(docx_file)
        assert "Test Document" in result["markdown"]

    def test_headings_converted(self, agent, docx_file):
        result = agent.process(docx_file)
        assert "# Section One" in result["markdown"]
        assert "## Subsection" in result["markdown"]

    def test_body_text_present(self, agent, docx_file):
        result = agent.process(docx_file)
        assert "This is a body paragraph." in result["markdown"]

    def test_table_converted_to_markdown(self, agent, docx_file):
        md = agent.process(docx_file)["markdown"]
        assert "| Name" in md
        assert "| alpha" in md

    def test_metadata_fields(self, agent, docx_file):
        meta = agent.process(docx_file)["metadata"]
        assert meta["file_type"] == "docx"
        assert meta["title"] == "Test Document"
        assert "source_file" in meta

    def test_no_headings_still_works(self, agent, tmp_path):
        f = _make_docx(tmp_path / "flat.docx", headings=False, tables=False)
        result = agent.process(f)
        assert "This is a body paragraph." in result["markdown"]


# ---------------------------------------------------------------------------
# PptxAgent tests
# ---------------------------------------------------------------------------

class TestPptxAgent:
    @pytest.fixture
    def agent(self, tmp_path):
        from src.agents.office_agent import PptxAgent
        return PptxAgent(image_agent=_mock_image_agent(), images_dir=tmp_path / "images")

    @pytest.fixture
    def pptx_file(self, tmp_path):
        return _make_pptx(tmp_path / "test.pptx")

    def test_returns_markdown_and_metadata(self, agent, pptx_file):
        result = agent.process(pptx_file)
        assert "markdown" in result
        assert "metadata" in result

    def test_slide_heading_present(self, agent, pptx_file):
        md = agent.process(pptx_file)["markdown"]
        assert "## Slide 1" in md

    def test_slide_title_present(self, agent, pptx_file):
        md = agent.process(pptx_file)["markdown"]
        assert "Slide Title" in md

    def test_bullet_content_present(self, agent, pptx_file):
        md = agent.process(pptx_file)["markdown"]
        assert "Bullet point content" in md

    def test_metadata_file_type(self, agent, pptx_file):
        meta = agent.process(pptx_file)["metadata"]
        assert meta["file_type"] == "pptx"
        assert "slide_count" in meta
        assert meta["slide_count"] >= 1

    def test_metadata_title(self, agent, pptx_file):
        meta = agent.process(pptx_file)["metadata"]
        assert meta["title"] == "Test Presentation"


# ---------------------------------------------------------------------------
# XlsxAgent tests
# ---------------------------------------------------------------------------

class TestXlsxAgent:
    @pytest.fixture
    def agent(self):
        from src.agents.office_agent import XlsxAgent
        return XlsxAgent()

    @pytest.fixture
    def xlsx_file(self, tmp_path):
        return _make_xlsx(tmp_path / "test.xlsx")

    def test_returns_markdown_and_metadata(self, agent, xlsx_file):
        result = agent.process(xlsx_file)
        assert "markdown" in result
        assert "metadata" in result

    def test_sheet_heading_present(self, agent, xlsx_file):
        md = agent.process(xlsx_file)["markdown"]
        assert "## Sheet: Data" in md

    def test_header_row_in_table(self, agent, xlsx_file):
        md = agent.process(xlsx_file)["markdown"]
        assert "| City" in md
        assert "| Population" in md

    def test_data_rows_present(self, agent, xlsx_file):
        md = agent.process(xlsx_file)["markdown"]
        assert "Tokyo" in md
        assert "Berlin" in md

    def test_metadata_file_type(self, agent, xlsx_file):
        meta = agent.process(xlsx_file)["metadata"]
        assert meta["file_type"] == "xlsx"
        assert "source_file" in meta

    def test_empty_sheet_handled(self, agent, tmp_path):
        import openpyxl
        wb = openpyxl.Workbook()
        wb.active.title = "Empty"
        f = tmp_path / "empty.xlsx"
        wb.save(str(f))
        md = agent.process(f)["markdown"]
        assert "Empty sheet" in md
